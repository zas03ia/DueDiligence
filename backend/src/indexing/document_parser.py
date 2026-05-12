import os
import uuid
from typing import Dict, Any, List, Optional
from pathlib import Path
import PyPDF2
from docx import Document
import openpyxl
from pptx import Presentation
import json
from src.models.enums import DocumentType
from src.utils.exceptions import IndexingError


class DocumentParser:
    """Multi-format document parser with metadata extraction"""
    
    def __init__(self):
        self.supported_types = {
            DocumentType.PDF: self._parse_pdf,
            DocumentType.DOCX: self._parse_docx,
            DocumentType.XLSX: self._parse_xlsx,
            DocumentType.PPTX: self._parse_pptx
        }
    
    def parse_document(self, file_path: str, document_type: DocumentType) -> Dict[str, Any]:
        """Parse document and extract text with metadata"""
        if not os.path.exists(file_path):
            raise IndexingError(f"File not found: {file_path}")
        
        parser_func = self.supported_types.get(document_type)
        if not parser_func:
            raise IndexingError(f"Unsupported document type: {document_type}")
        
        try:
            result = parser_func(file_path)
            result.update({
                "file_path": file_path,
                "file_type": document_type,
                "file_size": os.path.getsize(file_path)
            })
            return result
        except Exception as e:
            raise IndexingError(f"Failed to parse {document_type} document: {str(e)}")
    
    def _parse_pdf(self, file_path: str) -> Dict[str, Any]:
        """Parse PDF document with text and metadata"""
        text_content = []
        metadata = {}
        
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                # Extract metadata
                if pdf_reader.metadata:
                    metadata = {
                        "title": pdf_reader.metadata.get('/Title', ''),
                        "author": pdf_reader.metadata.get('/Author', ''),
                        "subject": pdf_reader.metadata.get('/Subject', ''),
                        "creator": pdf_reader.metadata.get('/Creator', ''),
                        "producer": pdf_reader.metadata.get('/Producer', ''),
                        "creation_date": str(pdf_reader.metadata.get('/CreationDate', '')),
                        "modification_date": str(pdf_reader.metadata.get('/ModDate', '')),
                        "pages": len(pdf_reader.pages)
                    }
                
                # Extract text from each page
                for page_num, page in enumerate(pdf_reader.pages):
                    try:
                        page_text = page.extract_text()
                        if page_text.strip():
                            text_content.append({
                                "page": page_num + 1,
                                "text": page_text.strip(),
                                "chunks": self._split_text_into_chunks(page_text.strip())
                            })
                    except Exception as e:
                        print(f"Warning: Failed to extract text from page {page_num + 1}: {e}")
                        continue
        
        except Exception as e:
            raise IndexingError(f"PDF parsing failed: {str(e)}")
        
        return {
            "text_content": text_content,
            "metadata": metadata,
            "total_text": " ".join([page["text"] for page in text_content])
        }
    
    def _parse_docx(self, file_path: str) -> Dict[str, Any]:
        """Parse DOCX document with text and metadata"""
        text_content = []
        metadata = {}
        
        try:
            doc = Document(file_path)
            
            # Extract metadata
            core_props = doc.core_properties
            metadata = {
                "title": core_props.title or '',
                "author": core_props.author or '',
                "subject": core_props.subject or '',
                "created": str(core_props.created) if core_props.created else '',
                "modified": str(core_props.modified) if core_props.modified else '',
                "keywords": core_props.keywords or '',
                "comments": core_props.comments or '',
                "paragraphs": len(doc.paragraphs)
            }
            
            # Extract text from paragraphs
            full_text = []
            for para_num, paragraph in enumerate(doc.paragraphs):
                if paragraph.text.strip():
                    para_text = paragraph.text.strip()
                    full_text.append(para_text)
            
            # Split into chunks
            combined_text = "\\n".join(full_text)
            text_content.append({
                "section": "main_content",
                "text": combined_text,
                "chunks": self._split_text_into_chunks(combined_text)
            })
            
        except Exception as e:
            raise IndexingError(f"DOCX parsing failed: {str(e)}")
        
        return {
            "text_content": text_content,
            "metadata": metadata,
            "total_text": combined_text
        }
    
    def _parse_xlsx(self, file_path: str) -> Dict[str, Any]:
        """Parse XLSX document with text and metadata"""
        text_content = []
        metadata = {}
        all_text = []
        
        try:
            workbook = openpyxl.load_workbook(file_path, read_only=True)
            
            # Extract metadata
            metadata = {
                "title": workbook.properties.title or '',
                "author": workbook.properties.creator or '',
                "subject": workbook.properties.subject or '',
                "created": str(workbook.properties.created) if workbook.properties.created else '',
                "modified": str(workbook.properties.modified) if workbook.properties.modified else '',
                "keywords": workbook.properties.keywords or '',
                "sheets": len(workbook.sheetnames)
            }
            
            # Extract text from each sheet
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_text = []
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = [str(cell) for cell in row if cell is not None and str(cell).strip()]
                    if row_text:
                        sheet_text.extend(row_text)
                        all_text.extend(row_text)
                
                if sheet_text:
                    combined_sheet_text = " ".join(sheet_text)
                    text_content.append({
                        "sheet": sheet_name,
                        "text": combined_sheet_text,
                        "chunks": self._split_text_into_chunks(combined_sheet_text)
                    })
            
            workbook.close()
            
        except Exception as e:
            raise IndexingError(f"XLSX parsing failed: {str(e)}")
        
        return {
            "text_content": text_content,
            "metadata": metadata,
            "total_text": " ".join(all_text)
        }
    
    def _parse_pptx(self, file_path: str) -> Dict[str, Any]:
        """Parse PPTX document with text and metadata"""
        text_content = []
        metadata = {}
        all_text = []
        
        try:
            prs = Presentation(file_path)
            
            # Extract metadata
            core_props = prs.core_properties
            metadata = {
                "title": core_props.title or '',
                "author": core_props.author or '',
                "subject": core_props.subject or '',
                "created": str(core_props.created) if core_props.created else '',
                "modified": str(core_props.modified) if core_props.modified else '',
                "keywords": core_props.keywords or '',
                "slides": len(prs.slides)
            }
            
            # Extract text from each slide
            for slide_num, slide in enumerate(prs.slides):
                slide_text = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                        all_text.append(shape.text.strip())
                
                if slide_text:
                    combined_slide_text = " ".join(slide_text)
                    text_content.append({
                        "slide": slide_num + 1,
                        "text": combined_slide_text,
                        "chunks": self._split_text_into_chunks(combined_slide_text)
                    })
        
        except Exception as e:
            raise IndexingError(f"PPTX parsing failed: {str(e)}")
        
        return {
            "text_content": text_content,
            "metadata": metadata,
            "total_text": " ".join(all_text)
        }
    
    def _split_text_into_chunks(self, text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
        """Split text into overlapping chunks"""
        if len(text) <= chunk_size:
            return [text]
        
        chunks = []
        start = 0
        
        while start < len(text):
            end = start + chunk_size
            
            # Try to break at word boundary
            if end < len(text):
                last_space = text.rfind(' ', start, end)
                if last_space > start:
                    end = last_space
            
            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)
            
            start = max(start + 1, end - overlap)
        
        return chunks
    
    def extract_bounding_boxes(self, file_path: str) -> List[Dict[str, Any]]:
        """Extract bounding box information for citations (placeholder for future OCR integration)"""
        # This is a placeholder for future implementation with OCR libraries
        # like pytesseract or pdfplumber for precise bounding box extraction
        return []
